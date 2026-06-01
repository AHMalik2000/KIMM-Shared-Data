from __future__ import annotations

import csv
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs"
DOCS = ROOT / "docs"
DELIVERABLES = ROOT / "deliverables"

LABELS = ["World", "Sports", "Business", "Sci/Tech"]

LSTM_CM = [
    [1717, 62, 66, 55],
    [14, 1859, 15, 12],
    [58, 13, 1652, 177],
    [58, 18, 121, 1703],
]

TRANSFORMER_CM = [
    [1731, 34, 64, 71],
    [31, 1820, 35, 14],
    [58, 8, 1634, 200],
    [46, 11, 93, 1750],
]


def read_csv(name: str) -> list[dict[str, str]]:
    with (OUT / name).open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


def fmt(value: str | float) -> str:
    return f"{float(value):.4f}"


def table_md(headers: list[str], rows: list[list[str]]) -> str:
    lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join(["---"] * len(headers)) + " |",
    ]
    for row in rows:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def failure_reason(row: dict[str, str]) -> str:
    text = row["text"]
    true_label = row["true_label_name"]
    lstm = row["lstm_pred_name"]
    transformer = row["transformer_pred_name"]
    if "Olympics" in text:
        return "Sports terms dominate the short text even though the dataset section is World."
    if "Intel" in text:
        return "A technology company story overlaps Business and Sci/Tech vocabulary."
    if "Wildfires" in text:
        return "Weather prediction language mixes science terms with world-event framing."
    if "Google IPO" in text:
        return "Finance and technology cues conflict; the IPO framing points toward Business."
    if "T. rex" in text:
        return "Very short science item gives limited context; the Transformer overweights business-like wording."
    if "IBM" in text:
        return "Company hiring story overlaps labor/business cues and computing-company Sci/Tech label."
    return f"Boundary case between {true_label}, {lstm}, and {transformer} news language."


def excerpt(text: str, limit: int = 96) -> str:
    clean = " ".join(text.split())
    if len(clean) <= limit:
        return clean
    return clean[: limit - 3].rstrip() + "..."


def truncation_note(row: dict[str, str]) -> str:
    text = row["text"]
    if "Olympics" in text:
        return "Possible; long live-report format may exceed 128 tokens."
    if "Intel" in text:
        return "Unlikely; company and product cues appear early."
    if "Wildfires" in text:
        return "Unlikely; wildfire/prediction cues appear early."
    if "Google IPO" in text:
        return "Unlikely; IPO cues appear early."
    if "T. rex" in text:
        return "No; short item fits within the limit."
    if "IBM" in text:
        return "Unlikely; hiring/company cues appear early."
    return "Uncertain."


def build_markdown() -> dict[str, str]:
    main = read_csv("main_results.csv")
    params = read_csv("parameter_counts.csv")
    dataset_ablation = read_csv("dataset_size_ablation.csv")
    seq_ablation = read_csv("sequence_length_ablation.csv")
    failures = read_csv("selected_misclassifications.csv")

    main_rows = [
        [r["model"], r["split"], fmt(r["loss"]), fmt(r["accuracy"]), fmt(r["macro_f1"])]
        for r in main
    ]
    param_rows = [[r["model"], f'{int(r["trainable_parameters"]):,}'] for r in params]
    hyperparam_rows = [
        ["Vocabulary size", "20,000", "20,000"],
        ["Embedding dimension", "128", "128"],
        ["Encoder depth", "1 bidirectional LSTM layer", "2 Transformer Encoder layers"],
        ["Hidden / feedforward size", "Hidden size 128 per direction", "Feedforward dimension 512"],
        ["Attention heads", "Not applicable", "4"],
        ["Pooling strategy", "Masked mean pooling", "Masked mean pooling"],
        ["Dropout", "0.2", "0.2"],
        ["Optimizer / learning rate", "Adam / 1e-3", "Adam / 1e-3"],
        ["Epochs / batch size", "6 / 64", "6 / 64"],
    ]
    baseline_accuracy = {
        r["model"]: float(r["valid_accuracy"])
        for r in dataset_ablation
        if float(r["train_fraction"]) == 0.25
    }
    ablation_rows = [
        [
            r["model"],
            f'{float(r["train_fraction"]):.0%}',
            fmt(r["valid_accuracy"]),
            fmt(r["valid_macro_f1"]),
            f'{(float(r["valid_accuracy"]) - baseline_accuracy[r["model"]]) * 100:+.2f} pp',
        ]
        for r in dataset_ablation
    ]
    seq_rows = [
        [r["model"], r["max_sequence_length"], fmt(r["valid_loss"]), fmt(r["valid_accuracy"]), fmt(r["valid_macro_f1"])]
        for r in seq_ablation
    ]
    failure_rows = [
        [
            str(i + 1),
            excerpt(r["text"]),
            r["true_label_name"],
            r["lstm_pred_name"],
            r["transformer_pred_name"],
            fmt(r["lstm_confidence"]),
            fmt(r["transformer_confidence"]),
            truncation_note(r),
            failure_reason(r),
        ]
        for i, r in enumerate(failures)
    ]

    one_pager = f"""# One-Page Project Plan And Result Summary

## Project Question

This project compares a from-scratch LSTM classifier and a from-scratch Transformer Encoder classifier on AG News text classification. The core question is how recurrent sequence modeling and self-attention differ in classification performance, convergence behavior, data efficiency, and failure patterns under a controlled PyTorch/Jupyter workflow.

## Dataset And Preprocessing

- Dataset source: Hugging Face `fancyzhx/ag_news`, loaded with `datasets.load_dataset("fancyzhx/ag_news")`.
- Split: 108,000 train, 12,000 validation from the official training split with `seed=42`, and 7,600 official test examples for final evaluation only.
- Labels: `World`, `Sports`, `Business`, `Sci/Tech`, encoded as integer classes `0, 1, 2, 3`.
- Pipeline: regex tokenization, vocabulary built from training data only, 20,000-token vocabulary, padding/truncation to 128 tokens for the main comparison, and padding masks for pooling/attention.

## Models And Training

- LSTM: embedding layer, one-layer bidirectional LSTM, masked mean pooling, dropout, and linear classifier; 2,825,220 trainable parameters.
- Transformer Encoder: embedding layer, learned positional embeddings, two encoder layers, four attention heads, feedforward dimension 512, masked mean pooling, dropout, and linear classifier; 2,973,444 trainable parameters.
- Shared training setup: Adam, learning rate `1e-3`, batch size 64, dropout 0.2, 6 epochs, seed 42, CPU runtime.

## Experiments And Final Findings

- Main test result: LSTM accuracy 0.9120 and macro F1 0.9118; Transformer Encoder accuracy 0.9125 and macro F1 0.9126. The Transformer Encoder is slightly ahead, but the gap is small.
- Required ablation: dataset size at 25%, 50%, and 100%. LSTM led at 25% and 50%, while Transformer Encoder led at 100%, supporting the hypothesis that LSTM is more stable with limited data and Transformer benefits more from full data.
- Optional extension: sequence length 128 vs 256. Longer 256-token inputs reduced validation metrics for both models, so 128 remains the better setting in this run.
- Failure analysis: six selected test errors cover both-wrong, LSTM-only wrong, and Transformer-only wrong cases. Most errors reflect Business/Sci-Tech overlap, news-section ambiguity, or short texts with misleading surface cues.
"""

    report = f"""# LSTM vs Transformer Encoder for AG News Classification

## Table of Figures

Figure 1: Training and validation loss curves for the main comparison

Figure 2: LSTM final test confusion matrix

Figure 3: Transformer Encoder final test confusion matrix

Table 1: Dataset split sizes

Table 2: Class distribution by split

Table 3: Trainable parameter counts

Table 4: Final model configuration and hyperparameters

Table 5: Main validation and test results

Table 6: Dataset-size ablation results

Table 7: Optional sequence-length ablation results

Table 8: Selected misclassified examples

## Introduction

Text classification is one of the standard ways to evaluate whether a neural network can turn raw language into useful categories. In this report, I compare an LSTM classifier and a Transformer Encoder classifier on the same AG News task. The research question is: when dataset source, preprocessing, training budget, and metrics are controlled, does the recurrent LSTM or the self-attention Transformer Encoder give better classification performance, data efficiency, and error behavior?

The comparison is interesting because the two models represent different assumptions about sequence data. The LSTM reads the text recurrently and carries information through a hidden state. The Transformer Encoder uses self-attention, allowing every token to interact with other tokens directly. In theory, the Transformer should have more flexible global context modeling, while the LSTM may be more stable when the amount of data is reduced.

## Dataset

AG News is a four-class news classification dataset. Each input is a short news text, and the model predicts one of four categories: `World`, `Sports`, `Business`, or `Sci/Tech`. The project uses Hugging Face `fancyzhx/ag_news` as the only dataset source. The official training split was divided into 108,000 training examples and 12,000 validation examples using `train_test_split(test_size=0.1, seed=42)`. The official test split contains 7,600 examples and was reserved for final evaluation only.

| Split | Samples | Use |
| --- | --- | --- |
| Train | 108,000 | Vocabulary construction and model fitting |
| Validation | 12,000 | Model comparison, ablation evaluation, and interpretation |
| Test | 7,600 | Final evaluation only |

| Split | World | Sports | Business | Sci/Tech |
| --- | --- | --- | --- | --- |
| Train | 26,991 | 26,966 | 27,100 | 26,943 |
| Validation | 3,009 | 3,034 | 2,900 | 3,057 |
| Test | 1,900 | 1,900 | 1,900 | 1,900 |

Labels were verified as integer class indices `0, 1, 2, 3`, matching the requirement for `torch.nn.CrossEntropyLoss`. For preprocessing, the text was tokenized with a reproducible regular-expression tokenizer. The vocabulary was built only from the training split to avoid vocabulary leakage. The vocabulary size was capped at 20,000 tokens, with reserved indices for padding and unknown tokens. Tokens outside this train-built vocabulary map to `<unk>`, and the 20,000-token cap was kept fixed for both architectures and all main ablations. The main comparison used a maximum sequence length of 128 tokens. Both models received the same padded and truncated token IDs, and padding masks were used during LSTM pooling and Transformer attention.

## Models

The first model is a bidirectional LSTM classifier. It uses an embedding layer, a one-layer bidirectional LSTM, masked mean pooling, dropout, and a final linear classifier. The second model is a Transformer Encoder classifier. It uses an embedding layer, learned positional embeddings, two Transformer Encoder layers, four attention heads, feedforward dimension 512, masked mean pooling, dropout, and a final linear classifier. Both models were trained from scratch, without pretrained embeddings or pretrained language models.

    {table_md(["Model", "Trainable parameters"], param_rows)}

    {table_md(["Setting", "LSTM", "Transformer Encoder"], hyperparam_rows)}

The parameter counts are reasonably comparable for this assignment. The Transformer Encoder has about 5.2% more trainable parameters than the LSTM, so I treat its small final metric advantage cautiously. Most parameters in both models come from the shared embedding table: `20,000 × 128 = 2,560,000` embedding parameters, or about 90.6% of the LSTM and 86.1% of the Transformer Encoder. The non-embedding encoder and classifier components differ by about 148,224 parameters, so the total-size mismatch is modest but still worth acknowledging when interpreting small metric gaps.

## Experiments

The main comparison trained both models for 6 epochs using Adam, learning rate `1e-3`, batch size 64, dropout 0.2, seed 42, and CPU execution. The evaluation metrics were accuracy and macro F1-score. Macro F1 was included because AG News has multiple classes and the assignment requires class-sensitive evaluation, even though the dataset is nearly balanced.

The required ablation changed one major factor: the fraction of training data. The tested fractions were 25%, 50%, and 100% of the training split. The hypothesis was that the LSTM would be more stable with less data, while the Transformer Encoder would benefit more from the full dataset because self-attention has greater capacity. The 100% ablation run is not numerically identical to the main-comparison run because the ablation helper shuffles the selected training indices even when the selected fraction is 100%; this changes the minibatch order while keeping the same examples, seed, and hyperparameters.

An optional extension compared maximum sequence lengths of 128 and 256 tokens. This tested whether longer context improved validation performance enough to justify the additional training cost. The notebook also suppresses the PyTorch nested-tensor prototype warning so that final outputs focus on project-relevant results rather than a harmless implementation warning.

## Results

{table_md(["Model", "Split", "Loss", "Accuracy", "Macro F1"], main_rows)}

Both models reached similar final performance. On the official test split, the LSTM achieved 0.9120 accuracy and 0.9118 macro F1. The Transformer Encoder achieved 0.9125 accuracy and 0.9126 macro F1. The Transformer Encoder therefore performed slightly better on the final test metrics, but the difference is small and should not be overinterpreted as a decisive architectural win.

The epoch-1 validation results show a convergence-speed difference: the LSTM reached 0.8957 validation accuracy after one epoch, while the Transformer Encoder reached 0.8835. This suggests that the LSTM reached a usable validation accuracy faster in the first stage of training, even though the Transformer Encoder slightly surpassed it by the final test evaluation.

The loss curves show that both models continued to reduce training loss while validation loss stopped improving and then increased by the final epoch. This indicates overfitting by epoch 6. The LSTM had a sharper training-validation loss gap by the final epoch, while the Transformer Encoder retained slightly lower validation and test loss. A reasonable next experiment would tune early stopping or regularization using the validation set.

The confusion matrices show that Sports was the easiest class for both models. The largest recurring confusion was between Business and Sci/Tech. For the LSTM, 177 Business items were predicted as Sci/Tech and 121 Sci/Tech items were predicted as Business. For the Transformer Encoder, 200 Business items were predicted as Sci/Tech and 93 Sci/Tech items were predicted as Business. This directional asymmetry suggests that the Transformer Encoder more aggressively predicts Sci/Tech for business stories with technology-company surface terms, while the LSTM more often maps Sci/Tech stories back toward Business.

Figures used for the final report:

![Training and validation loss curves](../outputs/main_loss_curves.png)

![LSTM final test confusion matrix](../outputs/lstm_confusion_matrix.png)

![Transformer Encoder final test confusion matrix](../outputs/transformer_confusion_matrix.png)

## Ablation Study

Hypothesis: the LSTM will be more stable with limited training data, while the Transformer Encoder will benefit more from the full dataset. The changed variable was training-set fraction. The controlled variables were dataset source, train/validation split, tokenizer, vocabulary, maximum sequence length, optimizer, learning rate, batch size, epoch count, metrics, and validation evaluation.

{table_md(["Model", "Training fraction", "Validation accuracy", "Validation macro F1", "Accuracy gain vs 25%"], ablation_rows)}

The ablation supports the hypothesis. At 25% training data, the LSTM slightly outperformed the Transformer Encoder. At 50%, the LSTM again had higher validation accuracy and macro F1. At 100%, the Transformer Encoder moved ahead. Quantitatively, the LSTM gained about +3.70 percentage points from 25% to 100% training data, while the Transformer Encoder gained about +4.47 percentage points. This provides numerical evidence that the Transformer Encoder benefited more from access to the full training split, while the LSTM was more data-efficient in lower-data settings.

The optional sequence-length ablation is shown below.

{table_md(["Model", "Max length", "Validation loss", "Validation accuracy", "Validation macro F1"], seq_rows)}

Increasing maximum sequence length from 128 to 256 did not improve validation performance. Both models performed worse at 256 tokens. Under the current training budget and architecture sizes, longer context added cost without improving classification quality.

## Failure Analysis

The table below gives six representative misclassified test examples. It includes cases missed by both models, cases missed only by the LSTM, and cases missed only by the Transformer Encoder.

{table_md(["#", "Input excerpt", "True label", "LSTM prediction", "Transformer prediction", "LSTM conf.", "Transformer conf.", "Truncation note", "Likely reason"], failure_rows)}

The selected errors show that the two models often fail on semantic boundary cases. Some examples contain surface words strongly associated with a different news section, such as an Olympics item labeled World but predicted as Sports by both models; this is also the example where truncation is most plausible because live-report text can exceed the 128-token limit. Other examples involve technology companies or IPOs, where Business and Sci/Tech labels overlap and the main disambiguating cues appear early enough that truncation is unlikely to be the primary cause. Example 3 is also informative for calibration: the LSTM's incorrect World prediction has only 0.4821 confidence, which makes it a borderline error rather than a strong commitment, while the Transformer Encoder is highly confident in the correct Sci/Tech label at 0.9785. Overall, the models fail similarly when the text has ambiguous section cues, but they also show different single-model errors: the LSTM misread some science/world-style items, while the Transformer misread some short Sci/Tech items as Business.

## Conclusion

Both from-scratch models performed well on AG News under the controlled setup, reaching about 91.2% final test accuracy. The Transformer Encoder achieved the best final test metrics, but only by a small margin. The LSTM converged faster in the first epoch and was stronger in lower-data ablation settings, while the Transformer Encoder benefited more from the full training set.

The required dataset-size ablation supports the original data-efficiency hypothesis. The optional sequence-length ablation showed that simply increasing maximum sequence length from 128 to 256 did not improve validation metrics. The most reasonable next experiment would be to tune early stopping or regularization, because both models show overfitting by the final epoch. A further limitation is that the ablations were run with a single seed; repeating the comparison across multiple seeds would make the data-efficiency conclusion more statistically robust.

## AI Usage Appendix

The complete notebook and generated output files are included in the project repository. The AI Usage Appendix is provided as a separate required appendix and also included in the combined report deliverable.

## References

AG News dataset via Hugging Face `fancyzhx/ag_news`.

PyTorch documentation for recurrent layers, Transformer Encoder layers, and `torch.nn.CrossEntropyLoss`.
"""
    report = report.replace("\n    |", "\n|")

    appendix = """# AI Usage Appendix

## Tools Used

Codex was used for project planning, code drafting support, debugging review, document review, notebook-result summarization, report editing, presentation drafting, and final deliverable packaging.

## What AI Was Used For

AI was used to translate the assignment requirements into a controlled experiment plan, check the scaffold against dataset-source and test-set rules, draft and review PyTorch notebook components, draft report and presentation text from verified notebook outputs, and identify concise interpretations of the metric tables, ablation results, loss curves, confusion matrices, and selected misclassifications. For code-level work, AI assistance was used to draft or review the model class structure, shared `run_epoch` training/evaluation utility, plotting functions, and confusion-matrix/failure-analysis table generation. All generated code was inspected and run cell-by-cell, and the masked mean pooling and Transformer padding-mask logic were manually checked against the attention masks.

## Incorrect, Incomplete, Or Misleading AI Outputs

Some early scaffold text still described outputs as pending after the notebook had already been run. That wording was corrected in the project README, report notes, and presentation notes. Some notebook markdown cells also contained interpretation placeholders after outputs were generated; those were replaced with result-aware interpretations. During code review, the team also checked that AI-assisted examples did not apply `softmax` before `CrossEntropyLoss`, did not normalize confusion matrices when raw error counts were needed, and did not use the official test split during model selection or ablation. The team also had to ensure that no claim treated the small Transformer Encoder metric lead as a large or decisive result.

## Team Modifications And Decisions

The team kept the project constrained to one dataset source, Hugging Face `fancyzhx/ag_news`, and required the official test split to be used only for final evaluation. The team selected dataset-size ablation as the required ablation and kept sequence length as an optional extension. The team decided to interpret the results conservatively: the Transformer Encoder was slightly better on the final test split, while the LSTM was more data-efficient in the 25% and 50% training-data settings. The team also added explicit checks for integer labels, train-only vocabulary construction, padding masks, and reproducible seeds rather than relying on generated code without verification.

## Effect Of AI Use

AI reduced the time needed to organize requirements, audit deliverables, draft code scaffolds, and write clear explanations. The final responsibility remains with the team: all notebook outputs, figures, tables, code, and claims must be reviewed and approved before submission.
"""

    slides_md = """# Final Presentation Slide Plan

Reference format: 4:3 academic seminar deck, white background, blue title rule, top-left slide numbering, bottom-left course/lab footer, and bottom-right page count.

## Slide 1 - Title

LSTM vs Transformer Encoder for AG News Text Classification.

## Slide 2 - Table of Contents

Introduction and Motivation; Dataset and Preprocessing; Model Development; Main Comparison Results; Ablation Study; Failure Analysis; Conclusion and Future Work.

## Slide 3 - Introduction

Claim: Both models are common sequence classifiers, but they process text differently.

## Slide 4 - Dataset

Claim: AG News provides a controlled four-class classification task.

## Slide 5 - Preprocessing

Claim: Both models use the same processed input pipeline.

## Slide 6 - Model Development

Claim: The two models are different but comparable in size.

## Slide 7 - Main Results

Claim: Both models perform similarly, with a slight Transformer Encoder advantage.

## Slide 8 - Training Behavior

Claim: Both models show overfitting by the final epoch.

## Slide 9 - Confusion Matrices

Claim: Business and Sci/Tech are the main source of class-level errors.

## Slide 10 - Ablation Study

Claim: LSTM is more stable with less data; Transformer benefits from the full dataset.

## Slide 11 - Conclusion

Claim: Architecture differences appear most clearly in data efficiency and failure patterns.

## Anticipated Q&A

Q1: Why is the Transformer Encoder better with more data?
A1: Its self-attention layers have higher capacity to model token interactions, so they appear to benefit more once the full training split is available.

Q2: Would early stopping change the conclusion?
A2: It might improve both models because validation loss rises by epoch 6, but the final test gap is already small, so the conservative conclusion would likely remain.

Q3: Why use macro F1 if AG News is balanced?
A3: Macro F1 is required by the assignment and still confirms that performance is not hiding a severe class-specific failure.

Q4: Why did 256 tokens hurt performance?
A4: Longer inputs add cost and noise under the same architecture and epoch budget; the key AG News cues usually appear early.
"""

    return {
        "one_pager": one_pager,
        "report": report,
        "appendix": appendix,
        "slides": slides_md,
    }


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_text(cell, text: str, bold: bool = False) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.bold = bold
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def style_doc(doc: Document, title: str) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)
    styles = doc.styles
    styles["Normal"].font.name = "Calibri"
    styles["Normal"].font.size = Pt(10.5)
    for style_name, size, color in [
        ("Heading 1", 15, RGBColor(46, 116, 181)),
        ("Heading 2", 12.5, RGBColor(46, 116, 181)),
        ("Heading 3", 11.5, RGBColor(31, 77, 120)),
    ]:
        style = styles[style_name]
        style.font.name = "Calibri"
        style.font.size = Pt(size)
        style.font.color.rgb = color
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(title)
    r.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(11, 37, 69)
    p.paragraph_format.space_after = Pt(4)


def add_para(doc: Document, text: str) -> None:
    p = doc.add_paragraph(text)
    p.paragraph_format.space_after = Pt(5)
    p.paragraph_format.line_spacing = 1.08


def add_bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(3)
        p.add_run(item)


def add_table(doc: Document, headers: list[str], rows: list[list[str]], widths: list[float] | None = None) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, header in enumerate(headers):
        set_cell_text(hdr[i], header, bold=True)
        set_cell_shading(hdr[i], "F2F4F7")
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            set_cell_text(cells[i], value)
    if widths:
        for row in table.rows:
            for i, width in enumerate(widths):
                row.cells[i].width = Inches(width)
    doc.add_paragraph()


def add_image(doc: Document, path: Path, caption: str, width: float = 6.1) -> None:
    if path.exists():
        doc.add_picture(str(path), width=Inches(width))
        p = doc.add_paragraph(caption)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(5)
        for run in p.runs:
            run.font.size = Pt(8.5)
            run.italic = True


def write_docx_files(markdown: dict[str, str]) -> None:
    DELIVERABLES.mkdir(exist_ok=True)
    main = read_csv("main_results.csv")
    params = read_csv("parameter_counts.csv")
    dataset_ablation = read_csv("dataset_size_ablation.csv")
    seq_ablation = read_csv("sequence_length_ablation.csv")
    failures = read_csv("selected_misclassifications.csv")

    # One-pager
    doc = Document()
    style_doc(doc, "One-Page Project Plan And Result Summary")
    add_para(doc, "Research question: compare a from-scratch LSTM classifier and Transformer Encoder classifier on AG News under one controlled PyTorch/Jupyter pipeline.")
    add_bullets(doc, [
        "Dataset: Hugging Face fancyzhx/ag_news only; 108,000 train, 12,000 validation, 7,600 official test examples.",
        "Preprocessing: train-only 20,000-token vocabulary, regex tokenizer, max length 128, padding masks for pooling and attention.",
        "Models: bidirectional LSTM with 2,825,220 parameters; Transformer Encoder with 2,973,444 parameters.",
        "Training: Adam, learning rate 1e-3, batch size 64, dropout 0.2, 6 epochs, seed 42.",
        "Final test result: Transformer Encoder slightly leads, 0.9125 accuracy and 0.9126 macro F1 vs LSTM 0.9120 and 0.9118.",
        "Ablation: LSTM leads at 25% and 50% data; Transformer Encoder leads at 100%, supporting the data-efficiency hypothesis.",
        "Failure analysis: six selected test errors show Business/Sci-Tech overlap and ambiguous news-section cues.",
    ])
    doc.save(DELIVERABLES / "ag_news_one_page_summary.docx")

    # Final report
    doc = Document()
    style_doc(doc, "LSTM vs Transformer Encoder for AG News Classification")
    sections = markdown["report"].split("\n## ")
    for index, section_text in enumerate(sections):
        if index == 0:
            continue
        lines = section_text.splitlines()
        doc.add_heading(lines[0], level=1)
        body = "\n".join(lines[1:]).strip()
        if lines[0].startswith("Model Development"):
            add_para(doc, "The LSTM uses an embedding layer, one-layer bidirectional LSTM, masked mean pooling, dropout, and a linear classifier. The Transformer Encoder uses embeddings, learned positional embeddings, two encoder layers, four heads, masked mean pooling, dropout, and a classifier.")
            add_table(doc, ["Model", "Trainable parameters"], [[r["model"], f'{int(r["trainable_parameters"]):,}'] for r in params], [3.0, 2.0])
            add_table(
                doc,
                ["Setting", "LSTM", "Transformer Encoder"],
                [
                    ["Vocabulary size", "20,000", "20,000"],
                    ["Embedding dimension", "128", "128"],
                    ["Encoder depth", "1 bidirectional LSTM layer", "2 Transformer Encoder layers"],
                    ["Hidden / feedforward", "128 per direction", "512 feedforward"],
                    ["Attention heads", "Not applicable", "4"],
                    ["Pooling", "Masked mean pooling", "Masked mean pooling"],
                    ["Dropout", "0.2", "0.2"],
                    ["Optimizer / LR", "Adam / 1e-3", "Adam / 1e-3"],
                    ["Epochs / batch", "6 / 64", "6 / 64"],
                ],
                [1.7, 2.2, 2.3],
            )
            add_para(doc, "Both models were trained from scratch with no pretrained embeddings or pretrained language models.")
        elif lines[0].startswith("Results and Interpretation"):
            add_table(doc, ["Model", "Split", "Loss", "Accuracy", "Macro F1"], [[r["model"], r["split"], fmt(r["loss"]), fmt(r["accuracy"]), fmt(r["macro_f1"])] for r in main], [2.2, 0.8, 0.8, 0.9, 0.9])
            add_para(doc, "Both models reached similar final performance. The Transformer Encoder was slightly better on final test metrics, but the gap was small.")
            add_image(doc, OUT / "main_loss_curves.png", "Figure 1. Training and validation loss curves.")
            add_image(doc, OUT / "lstm_confusion_matrix.png", "Figure 2. LSTM final test confusion matrix.", 5.3)
            add_image(doc, OUT / "transformer_confusion_matrix.png", "Figure 3. Transformer Encoder final test confusion matrix.", 5.3)
            add_para(doc, "Both models overfit by epoch 6. Sports was easiest for both models, and Business/Sci/Tech confusion was the main class-level error pattern.")
        elif lines[0].startswith("Ablation Study"):
            add_para(doc, "Hypothesis: LSTM is more stable with limited data, while Transformer Encoder benefits more from full data. The changed variable was training-set fraction; all other major settings were controlled.")
            add_table(doc, ["Model", "Train fraction", "Accuracy", "Macro F1"], [[r["model"], f'{float(r["train_fraction"]):.0%}', fmt(r["valid_accuracy"]), fmt(r["valid_macro_f1"])] for r in dataset_ablation], [2.4, 1.0, 1.0, 1.0])
            add_para(doc, "The ablation supports the hypothesis: LSTM led at 25% and 50%, while Transformer Encoder led at 100%.")
            add_table(doc, ["Model", "Max length", "Loss", "Accuracy", "Macro F1"], [[r["model"], r["max_sequence_length"], fmt(r["valid_loss"]), fmt(r["valid_accuracy"]), fmt(r["valid_macro_f1"])] for r in seq_ablation], [2.2, 0.9, 0.8, 0.9, 0.9])
            add_para(doc, "The optional sequence-length extension showed that 256-token inputs reduced validation performance for both models.")
        elif lines[0].startswith("Failure Analysis"):
            add_para(doc, "The selected examples cover both-wrong, LSTM-only wrong, and Transformer-only wrong cases.")
            rows = []
            for i, r in enumerate(failures):
                rows.append([
                    str(i + 1),
                    excerpt(r["text"], 72),
                    r["true_label_name"],
                    r["lstm_pred_name"],
                    r["transformer_pred_name"],
                    f'{float(r["lstm_confidence"]):.2f}',
                    f'{float(r["transformer_confidence"]):.2f}',
                    failure_reason(r),
                ])
            add_table(doc, ["#", "Input excerpt", "True", "LSTM", "Transformer", "LSTM c.", "Trans. c.", "Likely reason"], rows, [0.25, 1.45, 0.6, 0.7, 0.75, 0.45, 0.5, 2.0])
            add_para(doc, "Most remaining errors are semantic boundary cases involving class overlap or misleading surface cues, not label-format or preprocessing failures.")
        else:
            for para in body.split("\n\n"):
                clean = para.strip()
                if clean and not clean.startswith("|") and not clean.startswith("- `"):
                    add_para(doc, clean)
    doc.save(DELIVERABLES / "ag_news_final_report.docx")

    doc.add_section(WD_SECTION.NEW_PAGE)
    doc.add_heading("AI Usage Appendix", level=1)
    for para in markdown["appendix"].split("\n\n")[1:]:
        if para.startswith("## "):
            doc.add_heading(para[3:], level=2)
        else:
            add_para(doc, para)
    doc.save(DELIVERABLES / "ag_news_final_report_with_ai_appendix.docx")

    # Standalone appendix
    doc = Document()
    style_doc(doc, "AI Usage Appendix")
    for para in markdown["appendix"].split("\n\n")[1:]:
        if para.startswith("## "):
            doc.add_heading(para[3:], level=1)
        else:
            add_para(doc, para)
    doc.save(DELIVERABLES / "ag_news_ai_usage_appendix.docx")


def add_ppt_title(slide, kicker: str, title: str, slide_num: int | None = None) -> None:
    if slide_num is not None:
        num = slide.shapes.add_textbox(PptInches(0.08), PptInches(0.12), PptInches(0.25), PptInches(0.22))
        p = num.text_frame.paragraphs[0]
        p.text = str(slide_num)
        p.font.size = PptPt(13)
        p.font.bold = True
        p.font.color.rgb = PptRGBColor(31, 78, 121)
    box = slide.shapes.add_textbox(PptInches(0.28), PptInches(0.16), PptInches(9.25), PptInches(0.42))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = PptPt(15 if len(title) > 76 else 17)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(11, 37, 69)
    line = slide.shapes.add_shape(1, PptInches(0.28), PptInches(0.58), PptInches(9.25), PptInches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = PptRGBColor(31, 78, 121)
    line.line.color.rgb = PptRGBColor(31, 78, 121)


def add_ppt_bullets(slide, items: list[str], x=0.6, y=1.05, w=8.7, h=5.5, size=12) -> None:
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = PptPt(size)
        p.font.color.rgb = PptRGBColor(35, 35, 35)
        p.space_after = PptPt(5)


def add_metric(slide, x: float, y: float, value: str, label: str, color=(46, 116, 181)) -> None:
    shape = slide.shapes.add_shape(1, PptInches(x), PptInches(y), PptInches(2.45), PptInches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGBColor(242, 244, 247)
    shape.line.color.rgb = PptRGBColor(210, 216, 224)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = PptPt(24)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(*color)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = PptPt(11)
    p2.font.color.rgb = PptRGBColor(70, 70, 70)
    p2.alignment = PP_ALIGN.CENTER


def add_footer(slide, num: int, total: int = 11) -> None:
    left = slide.shapes.add_textbox(PptInches(0.12), PptInches(7.18), PptInches(3.5), PptInches(0.16))
    p0 = left.text_frame.paragraphs[0]
    p0.text = "Advanced Battery Manufacturing Systems"
    p0.font.size = PptPt(5.5)
    p0.font.color.rgb = PptRGBColor(31, 78, 121)
    box = slide.shapes.add_textbox(PptInches(9.2), PptInches(7.18), PptInches(0.55), PptInches(0.16))
    p = box.text_frame.paragraphs[0]
    p.text = f"{num} of {total}"
    p.font.size = PptPt(5.5)
    p.font.color.rgb = PptRGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.RIGHT


def write_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]
    slides = [
        ("Title", "LSTM vs Transformer Encoder for AG News Text Classification", []),
        ("Table of Contents", "Table of Contents", ["Introduction and Motivation", "Dataset and Preprocessing", "Model Development", "Main Comparison Results", "Ablation Study", "Failure Analysis", "Conclusion and Future Work"]),
        ("Introduction", "Why compare LSTM and Transformer Encoder models?", ["Both models are common sequence classifiers, but they process text differently.", "The LSTM carries information recurrently through hidden states.", "The Transformer Encoder uses self-attention to connect tokens directly.", "The goal is to compare behavior, not only final accuracy."]),
        ("Dataset", "AG News provides a controlled four-class text classification task.", ["Dataset source: Hugging Face fancyzhx/ag_news only.", "Classes: World, Sports, Business, and Sci/Tech.", "Split: 108,000 train, 12,000 validation, and 7,600 official test examples.", "The official test set is used only for final evaluation."]),
        ("Preprocessing", "Both models use the same processed input pipeline.", ["Regex tokenizer and train-only vocabulary construction.", "Vocabulary size capped at 20,000 tokens.", "Maximum sequence length is 128 for the main comparison.", "Padding masks are used during LSTM pooling and Transformer attention."]),
        ("Model Development", "The two models are different but comparable in size.", ["LSTM: embedding layer, one-layer bidirectional LSTM, masked mean pooling, dropout, classifier.", "Transformer Encoder: embeddings, learned positional encoding, two encoder layers, four heads, classifier.", "Trainable parameters: LSTM 2.83M; Transformer Encoder 2.97M.", "Both models are trained from scratch."]),
        ("Main Results", "Both models perform similarly, with a slight Transformer Encoder advantage.", ["LSTM test accuracy: 0.9120; macro F1: 0.9118.", "Transformer Encoder test accuracy: 0.9125; macro F1: 0.9126.", "The difference is small, so the conclusion should be cautious."]),
        ("Training Behavior", "Both models show overfitting by the final epoch.", ["Training loss continues to decrease across 6 epochs.", "Validation loss rises late in training.", "This suggests early stopping or stronger regularization as the next experiment."]),
        ("Confusion Matrices", "Business and Sci/Tech are the main source of class-level errors.", ["Sports is the easiest class for both models.", "Business and Sci/Tech overlap because many company stories include technology terms.", "The pattern is consistent with the selected failure examples."]),
        ("Ablation Study", "LSTM is more stable with less data; Transformer benefits from the full dataset.", ["25% data: LSTM 0.8719 vs Transformer 0.8699 validation accuracy.", "50% data: LSTM 0.8944 vs Transformer 0.8901.", "100% data: Transformer 0.9146 vs LSTM 0.9089.", "The result supports the data-efficiency hypothesis."]),
        ("Conclusion", "Architecture differences appear most clearly in data efficiency and failure patterns.", ["Transformer Encoder slightly wins the final test comparison.", "LSTM performs better in limited-data settings.", "Longer 256-token inputs did not improve validation metrics.", "Future work: early stopping, regularization, and class-level error reduction."]),
    ]
    image_map = {
        7: OUT / "main_loss_curves.png",
        8: OUT / "main_loss_curves.png",
        9: OUT / "transformer_confusion_matrix.png",
    }
    for idx, (kicker, title, bullets) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PptRGBColor(255, 255, 255)
        if idx == 1:
            title_box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.45), PptInches(8.4), PptInches(1.2))
            p = title_box.text_frame.paragraphs[0]
            p.text = title
            p.font.size = PptPt(24)
            p.font.bold = True
            p.font.color.rgb = PptRGBColor(0, 77, 153)
            p.alignment = PP_ALIGN.CENTER
            info = slide.shapes.add_textbox(PptInches(1.2), PptInches(3.25), PptInches(7.6), PptInches(1.4))
            tf = info.text_frame
            for i, text in enumerate(["Course: Introduction to Deep Learning", "AG News Classification Term Project", "Muhammad Abdullah Hassan Malik"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = PptPt(12 if i < 2 else 11)
                p.font.bold = i == 2
                p.alignment = PP_ALIGN.CENTER
            add_footer(slide, idx, len(slides))
        elif idx == 2:
            add_ppt_title(slide, kicker, title, idx)
            add_ppt_bullets(slide, bullets, x=0.75, y=1.05, w=8.3, h=4.8, size=13)
            add_footer(slide, idx, len(slides))
        elif idx == 6:
            add_ppt_title(slide, kicker, title, idx)
            add_metric(slide, 0.75, 1.35, "2.83M", "LSTM parameters")
            add_metric(slide, 3.45, 1.35, "2.97M", "Transformer parameters")
            add_ppt_bullets(slide, bullets, x=0.75, y=3.0, w=8.5, h=3.5, size=11.5)
            add_footer(slide, idx, len(slides))
        elif idx in image_map:
            img = image_map[idx]
            if img.exists():
                slide.shapes.add_picture(str(img), PptInches(5.35), PptInches(1.0), width=PptInches(4.15))
            add_ppt_title(slide, kicker, title, idx)
            add_ppt_bullets(slide, bullets, x=0.55, y=1.05, w=4.55, h=5.6, size=11)
            add_footer(slide, idx, len(slides))
        else:
            add_ppt_title(slide, kicker, title, idx)
            add_ppt_bullets(slide, bullets, x=0.75, y=1.05, w=8.5, h=5.7, size=12)
            add_footer(slide, idx, len(slides))
    prs.save(DELIVERABLES / "ag_news_lstm_transformer_presentation.pptx")


def main() -> None:
    markdown = build_markdown()
    DELIVERABLES.mkdir(exist_ok=True)
    (DOCS / "one_page_project_summary.md").write_text(markdown["one_pager"], encoding="utf-8")
    (DOCS / "final_report.md").write_text(markdown["report"], encoding="utf-8")
    (DOCS / "ai_usage_appendix_final.md").write_text(markdown["appendix"], encoding="utf-8")
    (DOCS / "final_presentation_script.md").write_text(markdown["slides"], encoding="utf-8")
    write_docx_files(markdown)
    write_pptx()


if __name__ == "__main__":
    main()
